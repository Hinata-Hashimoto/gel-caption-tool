"""PNG and PPTX export for gel caption tool."""
import io
import glob as _glob
from PIL import Image, ImageDraw, ImageFont
from annotation_model import AppState

# Margin widths added to the output image / slide
EXPORT_LEFT = 110   # pixels at original resolution
EXPORT_TOP = 140    # pixels at original resolution
BAND_R = 3          # px radius for band markers in PNG export


def _font(size: int):
    candidates = []
    # macOS: Hiragino (supports Japanese + Latin) — searched via glob
    candidates += sorted(_glob.glob("/System/Library/Fonts/ヒラギノ角ゴシック W3.ttc"))
    candidates += sorted(_glob.glob("/System/Library/Fonts/ヒラギノ*.ttc"))[:1]
    candidates += [
        "/System/Library/Fonts/Helvetica.ttc",
        "/System/Library/Fonts/SFNSText.ttf",
        "/Library/Fonts/Arial.ttf",
        # Windows Japanese
        "C:/Windows/Fonts/meiryo.ttc",
        "C:/Windows/Fonts/YuGothM.ttc",
        "C:/Windows/Fonts/msgothic.ttc",
        # Windows English
        "C:/Windows/Fonts/arial.ttf",
        # Linux
        "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
        "/usr/share/fonts/opentype/ipafont-gothic/ipag.ttf",
    ]
    for p in candidates:
        try:
            return ImageFont.truetype(p, size)
        except Exception:
            continue
    return ImageFont.load_default()


def _svg_escape(text: str) -> str:
    return (text
            .replace('&', '&amp;').replace('<', '&lt;')
            .replace('>', '&gt;').replace('"', '&quot;'))


def _draw_dashed_hline(draw, x0, x1, y, color=(80, 80, 80), width=1, dash=7, gap=4):
    x = x0
    while x < x1:
        draw.line([(x, y), (min(x + dash, x1), y)], fill=color, width=width)
        x += dash + gap


def _no_shadow(shape):
    """Remove drop shadow from a pptx shape by setting an empty effectLst."""
    from pptx.oxml.ns import qn
    from lxml import etree
    spPr = shape._element.spPr
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    etree.SubElement(spPr, qn('a:effectLst'))


def export_png(state: AppState, path: str, target_dpi: int = 300):
    img = state.current_image()
    if img is None:
        raise ValueError("No image loaded")

    img = img.convert("RGBA")
    iw, ih = img.size

    total_w = iw + EXPORT_LEFT
    total_h = ih + EXPORT_TOP

    # Auto-scale so output is at least 1 200 px wide (→ 4" at 300 DPI)
    MIN_W = 1200
    scale = max(1, -(-MIN_W // total_w))  # ceiling division; capped below
    scale = min(scale, 3)

    el = EXPORT_LEFT  * scale
    et = EXPORT_TOP   * scale
    sw = total_w      * scale
    sh = total_h      * scale
    tick = 30         * scale

    out = Image.new("RGBA", (sw, sh), (255, 255, 255, 255))
    img_up = img.resize((iw * scale, ih * scale), Image.LANCZOS) if scale > 1 else img
    out.paste(img_up, (el, et))
    draw = ImageDraw.Draw(out)

    font_l = _font(15 * scale)
    font_s = _font(max(8, int(state.sample_font_size * 1.3)) * scale)

    # Ladder annotations
    for ann in state.ladder_annotations:
        if ann.skipped:
            continue
        y = int(ann.y_img) * scale + et
        if state.ladder_line_style == "full":
            x_start, x_end = el, sw
            label_x = el - 6 * scale
        else:
            x_start, x_end = el - tick, el
            label_x = el - tick - 4 * scale
        _draw_dashed_hline(draw, x_start, x_end, y,
                           width=scale, dash=7 * scale, gap=4 * scale)
        size_str = str(int(ann.size)) if ann.size == int(ann.size) else str(ann.size)
        label = f"{size_str} {ann.unit}"
        draw.text((label_x, y), label, font=font_l, fill=(0, 0, 0), anchor="rm")

    # Sample annotations (rotated text above image)
    for ann in state.sample_annotations:
        xc = int(ann.x_img) * scale + el
        bbox = font_s.getbbox(ann.name)
        tw = bbox[2] - bbox[0]
        th = bbox[3] - bbox[1]
        tmp = Image.new("RGBA", (tw + 4, th + 4), (255, 255, 255, 0))
        ImageDraw.Draw(tmp).text((2, 2), ann.name, font=font_s, fill=(0, 0, 0))
        rot = tmp.rotate(90, expand=True)
        rw, rh = rot.size
        px = xc - rw // 2
        py = et - rh - 6 * scale
        if py < 0:
            py = 0
        out.paste(rot, (px, py), rot)

    # Region annotations
    rby = et // 5
    for reg in state.region_annotations:
        x1c = int(reg.x_start) * scale + el
        x2c = int(reg.x_end)   * scale + el
        draw.line([(x1c, rby), (x2c, rby)], fill=(0, 0, 0), width=max(1, scale))
        xmid = (x1c + x2c) // 2
        draw.text((xmid, rby - 2 * scale), reg.name,
                  font=font_s, fill=(0, 0, 0), anchor="mb")

    # Band markers
    r = BAND_R * scale
    for xi, yi in state.band_markers:
        xc = int(xi) * scale + el
        yc = int(yi) * scale + et
        draw.ellipse([(xc - r, yc - r), (xc + r, yc + r)], fill=(220, 0, 0))

    out.convert("RGB").save(path, format="PNG", dpi=(target_dpi, target_dpi))


def export_svg(state: AppState, path: str):
    """Export annotated gel as SVG (vector text/lines, raster image embedded)."""
    import base64, io as _io

    img = state.current_image()
    if img is None:
        raise ValueError("No image loaded")

    iw, ih = img.size
    total_w = iw + EXPORT_LEFT
    total_h = ih + EXPORT_TOP

    buf = _io.BytesIO()
    img.convert("RGB").save(buf, format="PNG")
    img_b64 = base64.b64encode(buf.getvalue()).decode()

    font_size   = 15
    sample_fs   = max(8, int(state.sample_font_size * 1.3))
    tick        = 30
    FONT_FAMILY = "Helvetica, Arial, sans-serif"

    def line_el(x1, y1, x2, y2, **kw):
        attrs = " ".join(f'{k.replace("_", "-")}="{v}"' for k, v in kw.items())
        return f'  <line x1="{x1}" y1="{y1}" x2="{x2}" y2="{y2}" {attrs}/>'

    svg = []
    svg.append('<?xml version="1.0" encoding="UTF-8"?>')
    svg.append(f'<svg xmlns="http://www.w3.org/2000/svg" '
               f'xmlns:xlink="http://www.w3.org/1999/xlink" '
               f'width="{total_w}" height="{total_h}" '
               f'viewBox="0 0 {total_w} {total_h}">')
    svg.append(f'  <rect width="{total_w}" height="{total_h}" fill="white"/>')
    svg.append(f'  <image x="{EXPORT_LEFT}" y="{EXPORT_TOP}" '
               f'width="{iw}" height="{ih}" '
               f'href="data:image/png;base64,{img_b64}"/>')

    # Ladder annotations
    rby = EXPORT_TOP // 5
    for ann in state.ladder_annotations:
        if ann.skipped:
            continue
        y = int(ann.y_img) + EXPORT_TOP
        if state.ladder_line_style == "full":
            x0, x1 = EXPORT_LEFT, total_w
            lx = EXPORT_LEFT - 6
        else:
            x0, x1 = EXPORT_LEFT - tick, EXPORT_LEFT
            lx = EXPORT_LEFT - tick - 4
        svg.append(line_el(x0, y, x1, y,
                           stroke="#505050", stroke_width="1",
                           stroke_dasharray="7,4"))
        size_str = str(int(ann.size)) if ann.size == int(ann.size) else str(ann.size)
        label = _svg_escape(f"{size_str} {ann.unit}")
        svg.append(f'  <text x="{lx}" y="{y}" text-anchor="end" '
                   f'dominant-baseline="middle" '
                   f'font-family="{FONT_FAMILY}" font-size="{font_size}" '
                   f'font-weight="bold">{label}</text>')

    # Sample annotations (rotated)
    for ann in state.sample_annotations:
        xc = int(ann.x_img) + EXPORT_LEFT
        yc = EXPORT_TOP * 3 // 4
        name = _svg_escape(ann.name)
        svg.append(f'  <text transform="rotate(-90,{xc},{yc})" '
                   f'x="{xc}" y="{yc}" '
                   f'text-anchor="middle" dominant-baseline="middle" '
                   f'font-family="{FONT_FAMILY}" '
                   f'font-size="{sample_fs}">{name}</text>')

    # Region annotations
    for reg in state.region_annotations:
        x1c = int(reg.x_start) + EXPORT_LEFT
        x2c = int(reg.x_end)   + EXPORT_LEFT
        svg.append(line_el(x1c, rby, x2c, rby, stroke="black", stroke_width="1"))
        xmid = (x1c + x2c) // 2
        name = _svg_escape(reg.name)
        svg.append(f'  <text x="{xmid}" y="{rby - 3}" '
                   f'text-anchor="middle" dominant-baseline="auto" '
                   f'font-family="{FONT_FAMILY}" '
                   f'font-size="{sample_fs}">{name}</text>')

    # Band markers
    for xi, yi in state.band_markers:
        xc = int(xi) + EXPORT_LEFT
        yc = int(yi) + EXPORT_TOP
        svg.append(f'  <circle cx="{xc}" cy="{yc}" r="{BAND_R}" fill="#dc0000"/>')

    svg.append('</svg>')

    with open(path, 'w', encoding='utf-8') as f:
        f.write('\n'.join(svg))


def export_pptx(state: AppState, path: str):
    """Export to PPTX using a fixed 25 cm slide width so coordinates are
    independent of the image's stored DPI metadata."""
    from pptx import Presentation
    from pptx.util import Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

    img = state.current_image()
    if img is None:
        raise ValueError("No image loaded")

    iw, ih = img.size

    def cm(x: float) -> int:
        return int(x * 360000)

    LEFT_CM  = 2.5
    TOP_CM   = 3.0
    SLIDE_W_CM = 25.0

    LEFT  = cm(LEFT_CM)
    TOP   = cm(TOP_CM)
    IMG_W = cm(SLIDE_W_CM) - LEFT
    scale = IMG_W / iw
    IMG_H = int(ih * scale)
    SLIDE_W = cm(SLIDE_W_CM)
    SLIDE_H = TOP + IMG_H

    def ix2e(x: float) -> int:
        return LEFT + int(x * scale)

    def iy2e(y: float) -> int:
        return TOP + int(y * scale)

    prs = Presentation()
    prs.slide_width  = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Gel image
    buf = io.BytesIO()
    img.convert("RGB").save(buf, format="PNG")
    buf.seek(0)
    slide.shapes.add_picture(buf, Emu(LEFT), Emu(TOP), Emu(IMG_W), Emu(IMG_H))

    # Ladder lines + label text boxes
    label_h = cm(0.55)
    LINE_H = 15000
    SHORT_TICK = cm(0.8)
    for ann in state.ladder_annotations:
        if ann.skipped:
            continue
        size_str = str(int(ann.size)) if ann.size == int(ann.size) else str(ann.size)
        label = f"{size_str} {ann.unit}"
        yc = iy2e(ann.y_img)

        if state.ladder_line_style == "full":
            line_x, line_w = LEFT, IMG_W
        else:
            line_x, line_w = LEFT - SHORT_TICK, SHORT_TICK
        rect = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Emu(line_x), Emu(yc - LINE_H // 2),
            Emu(line_w), Emu(LINE_H)
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(0x50, 0x50, 0x50)
        rect.line.fill.background()
        _no_shadow(rect)

        label_right = (LEFT - cm(0.15)) if state.ladder_line_style == "full" else (line_x - cm(0.1))
        txb = slide.shapes.add_textbox(
            Emu(0), Emu(yc - label_h // 2),
            Emu(label_right), Emu(label_h)
        )
        tf = txb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = label
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0, 0, 0)

    # Sample name text boxes
    box_w = TOP - cm(0.4)
    box_h = cm(0.55)
    center_y = TOP * 3 // 4
    for ann in state.sample_annotations:
        xc = ix2e(ann.x_img)
        txb = slide.shapes.add_textbox(
            Emu(xc - box_w // 2), Emu(center_y - box_h // 2),
            Emu(box_w), Emu(box_h)
        )
        txb.rotation = 270
        tf = txb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = ann.name
        run.font.size = Pt(state.sample_font_size)
        run.font.color.rgb = RGBColor(0, 0, 0)

    # Region annotations (black line + centered label)
    reg_bar_y   = TOP * 1 // 5
    reg_label_h = cm(0.55)
    reg_box_w   = cm(5.0)
    for reg in state.region_annotations:
        x1e = ix2e(reg.x_start)
        x2e = ix2e(reg.x_end)

        bar_rect = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Emu(x1e), Emu(reg_bar_y),
            Emu(x2e - x1e), Emu(15000)
        )
        bar_rect.fill.solid()
        bar_rect.fill.fore_color.rgb = RGBColor(0, 0, 0)
        bar_rect.line.fill.background()
        _no_shadow(bar_rect)

        xmid = (x1e + x2e) // 2
        txb = slide.shapes.add_textbox(
            Emu(xmid - reg_box_w // 2), Emu(reg_bar_y - reg_label_h - cm(0.05)),
            Emu(reg_box_w), Emu(reg_label_h)
        )
        tf = txb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = reg.name
        run.font.size = Pt(state.sample_font_size)
        run.font.color.rgb = RGBColor(0, 0, 0)

    # Band markers (red ovals, ~1 mm radius)
    r = cm(0.1)
    for xi, yi in state.band_markers:
        xc = ix2e(xi)
        yc = iy2e(yi)
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Emu(xc - r), Emu(yc - r),
            Emu(r * 2), Emu(r * 2)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xDC, 0x00, 0x00)
        shape.line.color.rgb = RGBColor(0xDC, 0x00, 0x00)
        _no_shadow(shape)

    prs.save(path)
